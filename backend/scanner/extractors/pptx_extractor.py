"""
PowerPoint文档提取器
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from .base import BaseExtractor
from models.schemas import FileInfo, DocumentMetrics


class PptxExtractor(BaseExtractor):
    """PowerPoint文档提取器"""
    
    def can_handle(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.pptx', '.ppt']
    
    def extract(self, file_path: Path, file_info: FileInfo) -> DocumentMetrics:
        """提取PPT文档指标"""
        metrics = DocumentMetrics()
        
        try:
            prs = Presentation(str(file_path))
            
            # 幻灯片数
            metrics.slide_count = len(prs.slides)
            metrics.page_count = len(prs.slides)
            
            total_chars = 0
            total_images = 0
            total_tables = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    # 文本框
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                total_chars += len(run.text)
                    
                    # 表格
                    if shape.has_table:
                        total_tables += 1
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    total_chars += len(cell.text)
                    
                    # 图片
                    if hasattr(shape, 'image'):
                        total_images += 1
            
            metrics.char_count = total_chars
            metrics.image_count = total_images
            metrics.table_count = total_tables
            
        except Exception as e:
            file_info.parse_success = False
            file_info.parse_error = str(e)
        
        return metrics
    
    def extract_text(self, file_path: Path) -> str:
        """提取文本内容"""
        try:
            prs = Presentation(str(file_path))
            texts = []
            
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    slide_texts.append(cell.text)
                if slide_texts:
                    texts.append('\n'.join(slide_texts))
            
            return '\n\n'.join(texts)
        except Exception:
            return ""
